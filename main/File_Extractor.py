import os
import json
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
import docx2txt
import pandas as pd
import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def ocr_image(image_path):
    try:
        return pytesseract.image_to_string(Image.open(image_path)).strip()
    except:
        return ""

def extract_text_from_csv(filepath):
    df = pd.read_csv(filepath)
    return df.to_string(index=False)

def extract_text_from_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath, data_only=True)
    text = ''
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            text += ' '.join(str(cell) for cell in row if cell is not None) + '\n'
    return text

def extract_text_and_images_from_pdf(filepath, img_dir):
    os.makedirs(img_dir, exist_ok=True)
    pdf_file = fitz.open(filepath)
    text = ''
    image_paths = []
    for page_num in range(len(pdf_file)):
        page = pdf_file.load_page(page_num)
        text += page.get_text()
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf_file.extract_image(xref)
            image_bytes = base_image['image']
            ext = base_image['ext']
            img_name = f"page{page_num+1}_img{img_index+1}.{ext}"
            img_path = os.path.join(img_dir, img_name)
            with open(img_path, 'wb') as f:
                f.write(image_bytes)
            image_paths.append(img_path)
    return text, image_paths

def extract_text_and_images_from_docx(filepath, img_dir):
    os.makedirs(img_dir, exist_ok=True)
    text = docx2txt.process(filepath, img_dir)
    image_files = [os.path.join(img_dir, f) for f in os.listdir(img_dir) if not f.startswith('~$')]
    return text, image_files

def extract_text_and_images_from_pptx(filepath, img_dir):
    prs = Presentation(filepath)
    os.makedirs(img_dir, exist_ok=True)
    text = ''
    image_files = []
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + '\n'
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                img_bytes = img.blob
                ext = img.ext
                img_name = f"slide{slide_num+1}_img{len(image_files)+1}.{ext}"
                img_path = os.path.join(img_dir, img_name)
                with open(img_path, 'wb') as f:
                    f.write(img_bytes)
                image_files.append(img_path)
    return text, image_files

def extract_all_text_and_ocr_native(folder_path):
    text_content = {}
    ocr_text_content = {}
    for filename in os.listdir(folder_path):
        if filename.startswith('~$'):
            continue
        filepath = os.path.join(folder_path, filename)
        ext = filename.lower().split('.')[-1]
        img_dir = os.path.join(folder_path, f"_imgs_{filename}")
        try:
            if ext == 'csv':
                text_content[filename] = extract_text_from_csv(filepath)
                ocr_text_content[filename] = ""
            elif ext == 'xlsx':
                text_content[filename] = extract_text_from_xlsx(filepath)
                ocr_text_content[filename] = ""
            elif ext == 'pdf':
                text, images = extract_text_and_images_from_pdf(filepath, img_dir)
                ocr_text = "".join([ocr_image(img) + "\n" for img in images])
                text_content[filename] = text
                ocr_text_content[filename] = ocr_text
            elif ext == 'docx':
                text, images = extract_text_and_images_from_docx(filepath, img_dir)
                ocr_text = "".join([ocr_image(img) + "\n" for img in images])
                text_content[filename] = text
                ocr_text_content[filename] = ocr_text
            elif ext == 'pptx':
                text, images = extract_text_and_images_from_pptx(filepath, img_dir)
                ocr_text = "".join([ocr_image(img) + "\n" for img in images])
                text_content[filename] = text
                ocr_text_content[filename] = ocr_text
            else:
                text_content[filename] = ""
                ocr_text_content[filename] = ""
        except Exception as e:
            print(f"Error processing {filename}: {e}")
            text_content[filename] = ""
            ocr_text_content[filename] = ""
    return text_content, ocr_text_content

# Usage example
course_folder = r"C:\Users\monat\Desktop\LearnWise\Course1"  # Set your path

text_json, ocr_json = extract_all_text_and_ocr_native(course_folder)

with open('text_extracted_native.json', 'w', encoding='utf-8') as f1, \
     open('text_extracted_from_images_native.json', 'w', encoding='utf-8') as f2:
    json.dump(text_json, f1, indent=2, ensure_ascii=False)
    json.dump(ocr_json, f2, indent=2, ensure_ascii=False)

print("Finished native text and OCR extraction.")
